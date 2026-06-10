#!/usr/bin/env python3
"""Generate the DATN PC Store presentation from the verified Markdown inputs.

The preferred renderer uses python-pptx. A dependency-free OOXML fallback is kept so
this repository can still generate the presentation in offline build environments.
No external assets, base64 payloads, or video are used.
"""

from __future__ import annotations

import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.sax.saxutils import escape

ROOT = Path(__file__).resolve().parents[1]
OUTLINE_PATH = ROOT / "output" / "presentation_outline.md"
SCRIPT_PATH = ROOT / "output" / "presentation_script.md"
PLACEHOLDER_PATH = ROOT / "output" / "slide_image_placeholders.md"
NOTES_PATH = ROOT / "output" / "presentation_notes.md"
PPTX_PATH = ROOT / "output" / "DATN_PC_Store_Gioi_Thieu.pptx"

BLUE = "0B4F9C"
DARK_BLUE = "06356B"
ACCENT = "1976D2"
LIGHT_BLUE = "EAF3FC"
PALE_BLUE = "F5F9FE"
TEXT = "20354A"
MUTED = "60758A"
WHITE = "FFFFFF"
BORDER = "B8D2EC"
GREEN = "16866B"

SLIDE_W = 13.333
SLIDE_H = 7.5
EMU = 914400


@dataclass
class SlideData:
    number: int
    title: str
    bullets: list[str]
    notes: str = ""
    placeholders: list[str] | None = None


def parse_sections(path: Path) -> dict[int, tuple[str, list[str]]]:
    if not path.exists():
        raise FileNotFoundError(f"Thiếu tệp đầu vào: {path}")
    sections: dict[int, tuple[str, list[str]]] = {}
    current: int | None = None
    title = ""
    body: list[str] = []
    for raw in path.read_text(encoding="utf-8").splitlines():
        match = re.match(r"^## Slide\s+(\d+)\s*\|\s*(.+?)\s*$", raw)
        if match:
            if current is not None:
                sections[current] = (title, body)
            current = int(match.group(1))
            title = match.group(2).strip()
            body = []
        elif current is not None and raw.strip():
            body.append(raw.strip())
    if current is not None:
        sections[current] = (title, body)
    return sections


def load_slides() -> list[SlideData]:
    outline = parse_sections(OUTLINE_PATH)
    scripts = parse_sections(SCRIPT_PATH)
    placeholders: dict[int, list[str]] = {}
    for line in PLACEHOLDER_PATH.read_text(encoding="utf-8").splitlines():
        match = re.match(r"^- Slide\s+(\d+)\s*\|\s*(.+)$", line.strip())
        if match:
            placeholders[int(match.group(1))] = [
                item.strip() for item in match.group(2).split("|") if item.strip()
            ]

    slides: list[SlideData] = []
    for number in sorted(outline):
        title, lines = outline[number]
        bullets = [line[2:].strip() for line in lines if line.startswith("- ")]
        note_lines = scripts.get(number, (title, []))[1]
        notes = " ".join(line for line in note_lines if not line.startswith("- ")).strip()
        slides.append(SlideData(number, title, bullets, notes, placeholders.get(number, [])))

    expected = list(range(1, 21))
    actual = [slide.number for slide in slides]
    if actual != expected:
        raise ValueError(f"Outline phải có đúng Slide 01 đến Slide 20; nhận được {actual}")
    return slides


def write_notes(slides: list[SlideData]) -> None:
    lines = [
        "# Presentation Notes - DATN PC Store",
        "",
        "> Lời thuyết trình được tạo từ `output/presentation_script.md` và đồng bộ theo thứ tự slide.",
        "",
    ]
    for slide in slides:
        lines.extend([f"## Slide {slide.number:02d} | {slide.title}", "", slide.notes, ""])
        if slide.placeholders:
            lines.append("**Ảnh cần bổ sung:** " + "; ".join(slide.placeholders))
            lines.append("")
    NOTES_PATH.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def add_textbox_pptx(slide, x, y, w, h, text, size, color, bold=False, align=None, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if fill or line:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill or WHITE)
        shape.line.color.rgb = RGBColor.from_string(line or fill or WHITE)
        box = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box = shape.text_frame
    box.clear()
    box.margin_left = box.margin_right = Inches(0.12)
    box.margin_top = box.margin_bottom = Inches(0.06)
    box.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = box.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)
    p.alignment = align or PP_ALIGN.LEFT
    return shape


def render_with_python_pptx(slides: list[SlideData]) -> None:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor.from_string(DARK_BLUE if data.number == 1 else WHITE)

        if data.number == 1:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(SLIDE_H))
            accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor.from_string(ACCENT); accent.line.fill.background()
            add_textbox_pptx(slide, 0.9, 1.25, 10.9, 1.25, data.title.upper(), 34, WHITE, True)
            add_textbox_pptx(slide, 0.95, 2.65, 8.8, 0.65, data.bullets[0], 22, "D5E9FF", True)
            add_textbox_pptx(slide, 0.95, 3.45, 5.9, 0.58, data.bullets[1], 18, WHITE)
            add_textbox_pptx(slide, 0.95, 4.15, 4.1, 0.52, data.bullets[2], 16, "9CCBFF", True)
            add_textbox_pptx(slide, 10.35, 4.35, 1.6, 1.6, "PC", 28, WHITE, True, PP_ALIGN.CENTER, ACCENT, ACCENT)
        else:
            add_textbox_pptx(slide, 0.55, 0.27, 0.58, 0.48, f"{data.number:02d}", 14, WHITE, True, PP_ALIGN.CENTER, ACCENT, ACCENT)
            add_textbox_pptx(slide, 1.25, 0.2, 11.15, 0.68, data.title, 25, DARK_BLUE, True)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.94), Inches(12.2), Inches(0.035))
            bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(BORDER); bar.line.fill.background()

            content_w = 6.25 if data.placeholders else 11.75
            y = 1.23
            bullet_h = min(0.82, 4.85 / max(len(data.bullets), 1))
            for index, bullet in enumerate(data.bullets):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(content_w), Inches(bullet_h - 0.08))
                card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string(PALE_BLUE if index % 2 == 0 else WHITE)
                card.line.color.rgb = RGBColor.from_string("D7E6F5")
                add_textbox_pptx(slide, 0.84, y + 0.05, 0.35, bullet_h - 0.18, "•", 20, ACCENT, True, PP_ALIGN.CENTER)
                add_textbox_pptx(slide, 1.22, y + 0.05, content_w - 0.72, bullet_h - 0.18, bullet, 15, TEXT)
                y += bullet_h

            if data.placeholders:
                ph_x, ph_w = 7.25, 5.45
                ph_h = 4.75 / len(data.placeholders)
                for index, label in enumerate(data.placeholders):
                    top = 1.25 + index * ph_h
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ph_x), Inches(top), Inches(ph_w), Inches(ph_h - 0.18))
                    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(LIGHT_BLUE)
                    shape.line.color.rgb = RGBColor.from_string(ACCENT); shape.line.width = Pt(1.5)
                    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                    add_textbox_pptx(slide, ph_x + 0.3, top + 0.28, ph_w - 0.6, ph_h - 0.75, label, 16, BLUE, True, PP_ALIGN.CENTER)
                    add_textbox_pptx(slide, ph_x + 0.55, top + ph_h - 0.58, ph_w - 1.1, 0.25, "Chèn ảnh chụp giao diện thực tế tại đây", 9, MUTED, False, PP_ALIGN.CENTER)
            else:
                add_textbox_pptx(slide, 9.72, 5.47, 2.55, 0.48, "ASP.NET CORE MVC", 12, WHITE, True, PP_ALIGN.CENTER, BLUE, BLUE)

        add_textbox_pptx(slide, 0.62, 7.04, 3.0, 0.25, "DATN PC Store", 9, WHITE if data.number == 1 else MUTED, True)
        add_textbox_pptx(slide, 11.85, 7.04, 0.8, 0.25, f"{data.number:02d}/20", 9, WHITE if data.number == 1 else MUTED, True, PP_ALIGN.RIGHT)

        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = data.notes
        except (AttributeError, NotImplementedError):
            pass

    prs.save(PPTX_PATH)


def emu(value: float) -> int:
    return int(value * EMU)


def run_xml(text: str, size: int, color: str, bold: bool = False) -> str:
    weight = ' b="1"' if bold else ""
    return (
        f'<a:r><a:rPr lang="vi-VN" sz="{size * 100}"{weight} dirty="0">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill><a:latin typeface="Aptos"/>'
        f'</a:rPr><a:t>{escape(text)}</a:t></a:r>'
    )


def shape_xml(shape_id: int, x: float, y: float, w: float, h: float, text: str = "", *,
              fill: str = WHITE, line: str = WHITE, size: int = 14, color: str = TEXT,
              bold: bool = False, radius: bool = False, align: str = "l", dashed: bool = False) -> str:
    geom = "roundRect" if radius else "rect"
    dash = '<a:prstDash val="dash"/>' if dashed else ""
    paragraphs = []
    for part in text.split("\n"):
        paragraphs.append(
            f'<a:p><a:pPr algn="{align}" marL="0" indent="0"/>{run_xml(part, size, color, bold)}'
            f'<a:endParaRPr lang="vi-VN" sz="{size * 100}"/></a:p>'
        )
    tx = "" if not text else (
        '<p:txBody><a:bodyPr wrap="square" lIns="109728" tIns="54864" rIns="109728" bIns="54864" anchor="ctr"/>'
        '<a:lstStyle/>' + "".join(paragraphs) + '</p:txBody>'
    )
    return (
        f'<p:sp><p:nvSpPr><p:cNvPr id="{shape_id}" name="Shape {shape_id}"/>'
        '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr>'
        f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
        f'<a:prstGeom prst="{geom}"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
        f'<a:ln w="12700"><a:solidFill><a:srgbClr val="{line}"/></a:solidFill>{dash}</a:ln>'
        f'</p:spPr>{tx}</p:sp>'
    )


def slide_xml(data: SlideData) -> str:
    shapes: list[str] = []
    sid = 2

    def add(*args, **kwargs):
        nonlocal sid
        shapes.append(shape_xml(sid, *args, **kwargs)); sid += 1

    bg = DARK_BLUE if data.number == 1 else WHITE
    add(0, 0, SLIDE_W, SLIDE_H, fill=bg, line=bg)
    if data.number == 1:
        add(0, 0, 0.22, SLIDE_H, fill=ACCENT, line=ACCENT)
        add(0.9, 1.25, 10.9, 1.25, data.title.upper(), fill=DARK_BLUE, line=DARK_BLUE, size=34, color=WHITE, bold=True)
        add(0.95, 2.65, 8.8, 0.65, data.bullets[0], fill=DARK_BLUE, line=DARK_BLUE, size=22, color="D5E9FF", bold=True)
        add(0.95, 3.45, 5.9, 0.58, data.bullets[1], fill=DARK_BLUE, line=DARK_BLUE, size=18, color=WHITE)
        add(0.95, 4.15, 4.1, 0.52, data.bullets[2], fill=DARK_BLUE, line=DARK_BLUE, size=16, color="9CCBFF", bold=True)
        add(10.35, 4.35, 1.6, 1.6, "PC", fill=ACCENT, line=ACCENT, size=28, color=WHITE, bold=True, radius=True, align="ctr")
    else:
        add(0.55, 0.27, 0.58, 0.48, f"{data.number:02d}", fill=ACCENT, line=ACCENT, size=14, color=WHITE, bold=True, radius=True, align="ctr")
        add(1.25, 0.2, 11.15, 0.68, data.title, fill=WHITE, line=WHITE, size=25, color=DARK_BLUE, bold=True)
        add(0.55, 0.94, 12.2, 0.035, fill=BORDER, line=BORDER)
        content_w = 6.25 if data.placeholders else 11.75
        y = 1.23
        bullet_h = min(0.82, 4.85 / max(len(data.bullets), 1))
        for index, bullet in enumerate(data.bullets):
            fill = PALE_BLUE if index % 2 == 0 else WHITE
            add(0.65, y, content_w, bullet_h - 0.08, fill=fill, line="D7E6F5", radius=True)
            add(0.84, y + 0.05, 0.35, bullet_h - 0.18, "•", fill=fill, line=fill, size=20, color=ACCENT, bold=True, align="ctr")
            add(1.22, y + 0.05, content_w - 0.72, bullet_h - 0.18, bullet, fill=fill, line=fill, size=15, color=TEXT)
            y += bullet_h
        if data.placeholders:
            ph_h = 4.75 / len(data.placeholders)
            for index, label in enumerate(data.placeholders):
                top = 1.25 + index * ph_h
                add(7.25, top, 5.45, ph_h - 0.18, label + "\n\nChèn ảnh chụp giao diện thực tế tại đây",
                    fill=LIGHT_BLUE, line=ACCENT, size=15, color=BLUE, bold=True, radius=True, align="ctr", dashed=True)
        else:
            add(9.72, 5.47, 2.55, 0.48, "ASP.NET CORE MVC", fill=BLUE, line=BLUE, size=12, color=WHITE, bold=True, radius=True, align="ctr")
    footer_color = WHITE if data.number == 1 else MUTED
    add(0.62, 7.04, 3.0, 0.25, "DATN PC Store", fill=bg, line=bg, size=9, color=footer_color, bold=True)
    add(11.85, 7.04, 0.8, 0.25, f"{data.number:02d}/20", fill=bg, line=bg, size=9, color=footer_color, bold=True, align="r")
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/>'
        '</p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
        + "".join(shapes) + '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
    )


def render_ooxml_fallback(slides: list[SlideData]) -> None:
    """Write a standards-compliant minimal PPTX when python-pptx is unavailable."""
    slide_count = len(slides)
    content_overrides = "".join(
        f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for i in range(1, slide_count + 1)
    )
    content_types = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
{content_overrides}</Types>'''
    slide_ids = "".join(f'<p:sldId id="{255+i}" r:id="rId{i+1}"/>' for i in range(1, slide_count + 1))
    presentation = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst><p:sldIdLst>{slide_ids}</p:sldIdLst>
<p:sldSz cx="12191695" cy="6858000" type="screen16x9"/><p:notesSz cx="6858000" cy="9144000"/></p:presentation>'''
    pres_rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>']
    pres_rels += [f'<Relationship Id="rId{i+1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>' for i in range(1, slide_count + 1)]
    pres_rels_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">' + "".join(pres_rels) + '</Relationships>'
    root_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/></Relationships>'''
    master = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:clrMap accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" bg1="lt1" bg2="lt2" folHlink="folHlink" hlink="hlink" tx1="dk1" tx2="dk2"/><p:sldLayoutIdLst><p:sldLayoutId id="1" r:id="rId1"/></p:sldLayoutIdLst><p:txStyles/></p:sldMaster>'''
    master_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>'''
    layout = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank"><p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sldLayout>'''
    layout_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>'''
    theme = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="DATN PC Store"><a:themeElements><a:clrScheme name="DATN"><a:dk1><a:srgbClr val="20354A"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="06356B"/></a:dk2><a:lt2><a:srgbClr val="F5F9FE"/></a:lt2><a:accent1><a:srgbClr val="1976D2"/></a:accent1><a:accent2><a:srgbClr val="16866B"/></a:accent2><a:accent3><a:srgbClr val="0B4F9C"/></a:accent3><a:accent4><a:srgbClr val="60758A"/></a:accent4><a:accent5><a:srgbClr val="B8D2EC"/></a:accent5><a:accent6><a:srgbClr val="EAF3FC"/></a:accent6><a:hlink><a:srgbClr val="1976D2"/></a:hlink><a:folHlink><a:srgbClr val="0B4F9C"/></a:folHlink></a:clrScheme><a:fontScheme name="Aptos"><a:majorFont><a:latin typeface="Aptos Display"/></a:majorFont><a:minorFont><a:latin typeface="Aptos"/></a:minorFont></a:fontScheme><a:fmtScheme name="DATN"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="12700"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme></a:themeElements></a:theme>'''
    core = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"><dc:title>DATN PC Store - Giới thiệu hệ thống</dc:title><dc:creator>DATN PC Store</dc:creator><dc:subject>Đồ án tốt nghiệp ASP.NET Core MVC</dc:subject><dc:description>Bài thuyết trình được tạo tự động từ source code đã xác minh.</dc:description></cp:coreProperties>'''
    app = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"><Application>Microsoft Office PowerPoint</Application><PresentationFormat>Widescreen</PresentationFormat><Slides>{slide_count}</Slides><Notes>0</Notes></Properties>'''
    slide_rel = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>'''

    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(PPTX_PATH, "w", zipfile.ZIP_DEFLATED) as package:
        files = {
            "[Content_Types].xml": content_types, "_rels/.rels": root_rels,
            "ppt/presentation.xml": presentation, "ppt/_rels/presentation.xml.rels": pres_rels_xml,
            "ppt/slideMasters/slideMaster1.xml": master,
            "ppt/slideMasters/_rels/slideMaster1.xml.rels": master_rels,
            "ppt/slideLayouts/slideLayout1.xml": layout,
            "ppt/slideLayouts/_rels/slideLayout1.xml.rels": layout_rels,
            "ppt/theme/theme1.xml": theme, "docProps/core.xml": core, "docProps/app.xml": app,
        }
        for name, payload in files.items():
            package.writestr(name, payload)
        for data in slides:
            package.writestr(f"ppt/slides/slide{data.number}.xml", slide_xml(data))
            package.writestr(f"ppt/slides/_rels/slide{data.number}.xml.rels", slide_rel)


def validate_output(slides: list[SlideData]) -> None:
    if not PPTX_PATH.exists() or PPTX_PATH.stat().st_size < 10_000:
        raise RuntimeError("File PowerPoint chưa được tạo hợp lệ.")
    with zipfile.ZipFile(PPTX_PATH) as package:
        package.testzip()
        generated = [name for name in package.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        if len(generated) != len(slides):
            raise RuntimeError(f"PPTX có {len(generated)} slide, dự kiến {len(slides)}.")
        xml_text = "\n".join(package.read(name).decode("utf-8") for name in generated)
        missing = [label for slide in slides for label in (slide.placeholders or []) if escape(label) not in xml_text]
        if missing:
            raise RuntimeError("Thiếu placeholder trong PPTX: " + ", ".join(missing))


def main() -> int:
    slides = load_slides()
    write_notes(slides)
    renderer = "python-pptx"
    try:
        import pptx  # noqa: F401
    except ImportError:
        renderer = "OOXML fallback (python-pptx chưa được cài trong môi trường ngoại tuyến)"
        render_ooxml_fallback(slides)
    else:
        render_with_python_pptx(slides)
    validate_output(slides)

    waiting = [(slide.number, slide.placeholders) for slide in slides if slide.placeholders]
    print(f"Đã tạo {len(slides)} slide bằng {renderer}.")
    print(f"PowerPoint: {PPTX_PATH}")
    print(f"Ghi chú: {NOTES_PATH}")
    print("Các slide chờ chèn ảnh:")
    for number, labels in waiting:
        print(f"- Slide {number:02d}: {'; '.join(labels or [])}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
